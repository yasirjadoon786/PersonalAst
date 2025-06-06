from .base_tool import Tool
from pptx import Presentation
from typing import Any, Optional
from abc import ABC, abstractmethod
from pydantic import BaseModel
import math
import requests
import json
import os

class SlideGenerationTool(Tool):
    name: str = "PowerPoint Slides Generator"
    description: str = "Creates a simple PowerPoint slides"
    action_type: str = "ppt_generate"
    input_format: str = """
{
  "title": "string (title of the presentation)",
  "filename": "string (output filename)",
  "slides": [
    {
      "title": "string (slide title)",
      "points": ["string (bullet point)", "string (bullet point)", ...]
    },
    ...
  ]
}"""
    
    def run(self, input_text: Any) -> str:
        """Generate a simple PowerPoint presentation"""
        try:
            # Parse input if it's a string
            if isinstance(input_text, str):
                try:
                    input_data = json.loads(input_text)
                except json.JSONDecodeError:
                    return "Error: Input must be valid JSON. Strictly follow the Input Format of PowerPoint Slides Generator tool"
            else:
                input_data = input_text
            
            # Get basic presentation info
            title = input_data.get('title', 'Untitled Presentation')
            slides = input_data.get('slides', [])
            filename = input_data.get('filename', f"{title.replace(' ', '_')}.pptx")
            if not filename.endswith('.pptx'):
                filename += '.pptx'
            
            # Create presentation
            prs = Presentation()
            
            # Add title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = title
            
            # Add content slides
            for slide_data in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_data.get('title', 'Untitled Slide')
                
                # Add bullet points to content placeholder
                content = slide.placeholders[1]
                tf = content.text_frame
                
                for point in slide_data.get('points', []):
                    p = tf.add_paragraph()
                    p.text = point
            
            # Save the presentation
            prs.save(filename)
            
            return f"Created presentation '{title}' with {len(slides)} content slides. Saved to {filename}"
            
        except Exception as e:
            return f"Error: {str(e)}"
